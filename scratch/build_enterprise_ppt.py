import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    BG_DARK = RGBColor(11, 15, 25)
    CARD_BG = RGBColor(15, 23, 42)
    AMBER = RGBColor(245, 158, 11)
    EMERALD = RGBColor(16, 185, 129)
    SKY = RGBColor(56, 189, 248)
    PURPLE = RGBColor(168, 85, 247)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    ROSE = RGBColor(244, 63, 94)

    screenshots_dir = os.path.join(os.path.dirname(__file__), 'screenshots')

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_slide_header(slide, category, title):
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.9))
        tf = tx.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = category.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = AMBER

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

    def add_screen_slide(category, title, img_filename, aws_configs, key_features):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide)
        add_slide_header(slide, category, title)

        # Embedded Live Browser Screenshot (Left Side)
        img_path = os.path.join(screenshots_dir, img_filename)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.4), width=Inches(7.2))

        # AWS Config & Technical Details Panel (Right Side)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.4), Inches(4.7), Inches(5.6))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = AMBER
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "⚡ AWS Configuration Specs:"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = AMBER

        for cfg_title, cfg_val in aws_configs:
            p_c = tf.add_paragraph()
            r1 = p_c.add_run()
            r1.text = f"• {cfg_title}: "
            r1.font.bold = True
            r1.font.size = Pt(11)
            r1.font.color.rgb = SKY
            r2 = p_c.add_run()
            r2.text = cfg_val
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = WHITE

        p_space = tf.add_paragraph()
        p_space.text = " "
        p_space.font.size = Pt(4)

        p_f = tf.add_paragraph()
        p_f.text = "🛠️ Functional & ML Capabilities:"
        p_f.font.size = Pt(13)
        p_f.font.bold = True
        p_f.font.color.rgb = EMERALD

        for feat in key_features:
            p_feat = tf.add_paragraph()
            r1 = p_feat.add_run()
            r1.text = f"✓ {feat}"
            r1.font.size = Pt(10.5)
            r1.font.color.rgb = GRAY

    # SLIDE 1: Title Slide
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    box1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf1 = box1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "AWS CLOUD QUEST & ENTERPRISE AI ARCHITECTURE DECK"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AMBER

    p = tf1.add_paragraph()
    p.text = "PayComprehend AI: Nordic Financial Risk Intelligence"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p = tf1.add_paragraph()
    p.text = "Serverless Intelligent Document Processing (IDP) • Amazon Textract + Comprehend + Bedrock GenAI + Nordic Risk Engine"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY

    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.0), Inches(11.333), Inches(2.6))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = AMBER
    card1.line.width = Pt(1)

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    items = [
        ("⚡ 99.1% Operational Cost Reduction:", " Manual invoice entry ($18.00/doc) reduced to $0.015/doc via AWS Serverless OCR & NLP."),
        ("🇸🇪 Nordic Financial Riskklass Engine:", " Automatic calculation of Riskklass 1-5, Swedish Org.nr validation, F-skatt status & Bankgiro audit."),
        ("🧠 Amazon Bedrock (Claude 3.5 Sonnet):", " Deep financial mismatch explanations & insolvency probability scoring."),
        ("🛡️ Heavy AWS Services Safeguard:", " Itemized cost confirmation modal preventing unexpected AWS cloud spend."),
        ("☁️ CloudFormation 1-Click Teardown:", " 100% stack cleanup guarantee including S3 buckets, Lambda & CloudWatch Log Groups.")
    ]
    for idx, (head, body) in enumerate(items):
        p = tf_c1.paragraphs[0] if idx == 0 else tf_c1.add_paragraph()
        run1 = p.add_run()
        run1.text = head
        run1.font.bold = True
        run1.font.color.rgb = AMBER if idx == 0 else (EMERALD if idx == 1 else SKY)
        run1.font.size = Pt(11.5)
        run2 = p.add_run()
        run2.text = body
        run2.font.color.rgb = WHITE
        run2.font.size = Pt(11.5)

    # SLIDE 2: Standby Hero Screen
    add_screen_slide(
        "Screen 01 - Platform Standby State",
        "Interactive AWS AI Command Center Hero Screen",
        "10_standby_hero_screen.png",
        [
            ("AWS Region", "us-east-1 (N. Virginia)"),
            ("S3 Ingestion Bucket", "payment-docs-ai-dev-69753455"),
            ("Lambda Function", "PaymentDocProcessorFunction-dev"),
            ("Telemetry Status", "Operational (Latency < 1.8s)")
        ],
        [
            "Renders when all top header tabs are collapsed.",
            "Features central animated glowing AWS AI orb.",
            "Includes 4 quick-action launch module cards.",
            "Real-time telemetry bar showing active pipeline status."
        ]
    )

    # SLIDE 3: Executive Dashboard
    add_screen_slide(
        "Screen 02 - Executive Dashboard Module",
        "Real-Time Financial Spending Analytics & OCR Confidence",
        "01_dashboard_screen.png",
        [
            ("Analytics Source", "S3 JSON extraction records"),
            ("OCR Engine", "Amazon Textract AnalyzeExpense"),
            ("NLP Engine", "Amazon Comprehend NER & Sentiment"),
            ("Metric Cards", "Total Spend, File Count, Top Vendor, Confidence, Risk Alerts")
        ],
        [
            "Real-time expenditure aggregation across uploaded invoices.",
            "Top Spend Vendor identification & expense breakdown.",
            "Interactive expenditure timeline & vendor spend bar charts.",
            "Comprehend NLP highlights showing sentiment & risk alerts."
        ]
    )

    # SLIDE 4: Upload Station
    add_screen_slide(
        "Screen 03 - Upload Station & S3 Direct Ingestion",
        "Direct AWS S3 Upload & Collapsible 4-Step Pipeline Flow",
        "02_upload_screen.png",
        [
            ("S3 CORS Policy", "AllowedOrigins: [*], AllowedMethods: [PUT, POST]"),
            ("SDK Ingestion", "@aws-sdk/client-s3 PutObjectCommand"),
            ("Supported Formats", ".PDF, .PNG, .JPG, .DOC, .DOCX"),
            ("Lambda Trigger", "S3 ObjectCreated notification payload")
        ],
        [
            "Drag-and-drop file uploader supporting multi-format files.",
            "Direct S3 bucket upload bypassing server bottlenecks.",
            "Collapsible 4-step pipeline workflow panel (closed by default).",
            "Automatic fallback for universal file parsing."
        ]
    )

    # SLIDE 5: Document Explorer
    add_screen_slide(
        "Screen 04 - Document Explorer Module",
        "Nordic Financial Riskklass 1-5 Badges & Table Filtering",
        "03_explorer_screen.png",
        [
            ("Risk Klass Algorithm", "Nordic Swedish Credit Risk Matrix (0-100)"),
            ("Org.nr Validator", "Modulus-10 Luhn checksum check"),
            ("Tax Compliance", "Godkänd för F-skatt status check"),
            ("Payment Method Audit", "Bankgiro / Plusgiro matching")
        ],
        [
            "Interactive search across vendor, file name & invoice numbers.",
            "Filter by category, processing status & spend sorting.",
            "Displays Nordic Riskklass badges (Guld 5 to Röd 1).",
            "Exact currency formatting (£, €, $) per extracted document."
        ]
    )

    # SLIDE 6: Document Detail Modal
    add_screen_slide(
        "Screen 05 - Deep Document Audit & Bedrock GenAI",
        "Amazon Bedrock (Claude 3.5 Sonnet) & Rekognition Audit",
        "04_document_detail_modal.png",
        [
            ("GenAI LLM Model", "Anthropic Claude 3.5 Sonnet via Bedrock"),
            ("OCR Key-Value", "Textract KeyValuePairs & LineItemGroups"),
            ("Visual Audit", "Amazon Rekognition Stamp & Signature Detection"),
            ("Confidence Metric", "Weighted OCR + NLP confidence score")
        ],
        [
            "Renders side-by-side original invoice image & parsed data.",
            "Displays Bedrock GenAI natural language risk reasoning.",
            "Swedish Org.nr Modulus-10 & F-skatt status validation box.",
            "1-Click JSON data export matching AWS Studio Console."
        ]
    )

    # SLIDE 7: Architecture & ROI View
    add_screen_slide(
        "Screen 06 - Architecture & ROI Specification",
        "AWS End-to-End System Design & Resource Catalog",
        "05_architecture_screen.png",
        [
            ("Architecture Spec", "Event-driven Serverless Cloud Formation"),
            ("Resource Catalog", "S3, API Gateway, Lambda, Textract, Comprehend, Bedrock"),
            ("ROI Metric", "$0.015 / invoice (99.1% cost reduction)"),
            ("Processing Latency", "< 1.8s end-to-end execution speed")
        ],
        [
            "Embedded high-resolution 16:9 AWS Architecture Diagram.",
            "Collapsible Architecture Diagram panel (closed by default).",
            "Collapsible AWS Resource Catalog panel (closed by default).",
            "Itemized performance metrics & CloudFormation ARNs."
        ]
    )

    # SLIDE 8: Heavy Services Safeguard Modal
    add_screen_slide(
        "Screen 07 - Heavy Services Cost Safeguard Modal",
        "Itemized Pricing Transparency & User Confirmation",
        "06_paid_services_modal.png",
        [
            ("Bedrock Pricing", "~$0.003 / 1k tokens (~$0.005/doc)"),
            ("SageMaker Pricing", "~$0.00002 / compute sec (~$0.001/doc)"),
            ("Fraud Detector Pricing", "~$0.01 / fraud prediction"),
            ("Kendra Pricing", "~$0.81 / hour (~$590/month dev tier)")
        ],
        [
            "Protects AWS account against unexpected heavy cloud charges.",
            "Provides itemized per-document & per-hour cost breakdown.",
            "Requires explicit user disclaimer checkbox before enabling.",
            "Free standard tier ($0.00 base cost) remains 100% active."
        ]
    )

    # SLIDE 9: CloudFormation Visualizer
    add_screen_slide(
        "Screen 08 - CloudFormation Teardown Visualizer",
        "1-Click Infrastructure Deployment & Complete Teardown",
        "07_cloudformation_screen.png",
        [
            ("CloudFormation Template", "aws-backend/template.yaml"),
            ("Stack Name", "payment-ai-stack"),
            ("Log Group Resource", "PaymentDocProcessorLogGroup (AWS::Logs::LogGroup)"),
            ("S3 Teardown Hook", "Pre-deletion DeleteObjectsCommand execution")
        ],
        [
            "Interactive step-by-step stack creation visualizer.",
            "Glassmorphism Non-Empty S3 Bucket Warning Modal Popup.",
            "1-Click Teardown cleans up S3, Lambda, API Gateway & Log Group.",
            "Guarantees 100% deletion of resources for $0.00 leftover cost."
        ]
    )

    # SLIDE 10: AWS Backend & Lambda Code
    add_screen_slide(
        "Screen 09 - AWS Backend Code & Lambda Guide",
        "Syntax-Highlighted CloudFormation & Python 3.12 Lambda",
        "08_lambda_code_screen.png",
        [
            ("Lambda Runtime", "Python 3.12 with Boto3 SDK"),
            ("Textract Method", "textract.analyze_expense()"),
            ("Comprehend Method", "comprehend.detect_entities()"),
            ("CloudFormation Spec", "8 Managed Stack Resources in YAML")
        ],
        [
            "Panel 1: Collapsible CloudFormation Spec (template.yaml).",
            "Panel 2: Collapsible Python 3.12 Lambda Code (lambda_function.py).",
            "1-Click copy to clipboard for both YAML & Python source code.",
            "Step-by-step Boto3 SDK execution flow pills."
        ]
    )

    # SLIDE 11: AWS Config Modal
    add_screen_slide(
        "Screen 10 - AWS Config & Credentials Settings",
        "AWS Region, Access Keys & S3 Bucket Settings",
        "09_aws_config_modal.png",
        [
            ("Access Key ID", "AKIA... (User AWS Credentials)"),
            ("Secret Access Key", "K8sX... (Encrypted Local Storage)"),
            ("Target Region", "us-east-1 (N. Virginia)"),
            ("Target S3 Bucket", "payment-ai-stack-paymentdocumentbucket-yjxtkpxjyr63")
        ],
        [
            "Configures live AWS Boto3 / JS SDK authentication credentials.",
            "Live mode vs Lab Simulator toggle switch.",
            "Stores credentials safely in local browser storage.",
            "Allows custom S3 bucket & Lambda function targeting."
        ]
    )

    # SLIDE 12: How to Run Guide
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_slide_header(s12, "Developer Setup & Execution", "How to Run PayComprehend AI Locally")

    steps_code = [
        ("Step 1: Open Terminal in Project Directory", "cd \"C:\\Users\\CT0514\\OneDrive - Suomen Asiakastieto Oy\\sandbox\\workspaces\\comprehand\""),
        ("Step 2: Install Node Dependencies", "npm install"),
        ("Step 3: Launch Local Vite Development Server", "npm run dev"),
        ("Step 4: Open Application in Web Browser", "Navigate to http://localhost:3000/ in Chrome / Edge / Firefox")
    ]

    for idx, (step_title, code_str) in enumerate(steps_code):
        y = Inches(1.5 + idx * 1.35)
        c = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(1.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = AMBER
        c.line.width = Pt(1)

        tf = c.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = AMBER

        p = tf.add_paragraph()
        p.text = f"$ {code_str}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = EMERALD

    # Save PPTX
    out_dir = r"C:\Users\CT0514\OneDrive - Suomen Asiakastieto Oy\sandbox\workspaces\comprehand"
    out_path = os.path.join(out_dir, "CloudQuest_PayComprehend_AI_Presentation.pptx")
    prs.save(out_path)
    print(f"COMPLETE: 12-Slide Presentation saved successfully with screenshots & AWS Config details at: {out_path}")

if __name__ == "__main__":
    build_presentation()
